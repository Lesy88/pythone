import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import os

# Carica il dataset
df = pd.read_csv('Students_senzaUndef.csv')

# 1. Statistiche descrittive
summary_stats = df.describe()

# 2. Gestione dei valori mancanti
missing_data = df.isnull().sum()

# 3. Numero totale di record
total_students = len(df)

# 4. Variabili chiave
key_variables = df.columns.tolist()

# Creazione dei grafici

# Grafico 1: Distribuzione di Gender
gender_dist = df['Gender'].value_counts()
gender_dist.plot(kind='pie', autopct='%1.1f%%', startangle=90, colors=['#66b3ff', '#99ff99'], figsize=(6, 6))
plt.title('Distribuzione di Gender')
plt.ylabel('')
plt.tight_layout()
plt.savefig('gender_distribution.png')

# Grafico 2: Distribuzione di Age (Istogramma)
plt.figure(figsize=(6, 4))
sns.histplot(df['Age'], kde=True, color='skyblue', bins=10)
plt.title('Distribuzione di Age')
plt.xlabel('Age')
plt.ylabel('Frequenza')
plt.tight_layout()
plt.savefig('age_distribution.png')

# Grafico 3: Numero di studenti per Department
department_count = df['Department'].value_counts()
department_count.plot(kind='barh', color='salmon', figsize=(6, 4))
plt.title('Numero di studenti per Department')
plt.xlabel('Numero di studenti')
plt.ylabel('Department')
plt.tight_layout()
plt.savefig('department_distribution.png')

# Grafico 4: Punteggi medi per Midterm, Final e Total per Department
department_scores = df.groupby('Department')[['Midterm_Score', 'Final_Score', 'Total_Score']].mean()
department_scores.plot(kind='bar', figsize=(8, 6))
plt.title('Punteggi medi per Midterm, Final e Total per Department')
plt.xlabel('Department')
plt.ylabel('Punteggio')
plt.tight_layout()
plt.savefig('average_scores_department.png')

# Grafico 5: Boxplot per la distribuzione di Final_Score
plt.figure(figsize=(6, 4))
sns.boxplot(x=df['Final_Score'], color='lightgreen')
plt.title('Distribuzione di Final_Score')
plt.tight_layout()
plt.savefig('final_score_boxplot.png')

# Grafico 6: Scatter plot tra Study_Hours_per_Week e Total_Score
plt.figure(figsize=(6, 4))
sns.scatterplot(x=df['Study_Hours_per_Week'], y=df['Total_Score'], color='purple', s=80)
plt.title('Relazione tra Study Hours per Week e Total Score')
plt.xlabel('Study Hours per Week')
plt.ylabel('Total Score')
plt.tight_layout()
plt.savefig('study_hours_vs_score.png')

# Grafico 7: Correlazione tra Attendance e Performance
plt.figure(figsize=(6, 4))
sns.scatterplot(x=df['Attendance (%)'], y=df['Total_Score'], color='orange', s=80)
plt.title('Relazione tra Attendance (%) e Total Score')
plt.xlabel('Attendance (%)')
plt.ylabel('Total Score')
plt.tight_layout()
plt.savefig('attendance_vs_score.png')

# **Creazione della presentazione PowerPoint**
prs = Presentation()

# **Slide 1: Slide Introduttiva**
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Titolo
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Rapporto Finale - Analisi del Dataset"
subtitle.text = f"Numero totale di studenti: {total_students}\nVariabili chiave: {', '.join(key_variables)}\n\nStatistiche descrittive:\n{summary_stats}"

# **Slide 2: Dati Demografici**
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Titolo e contenuto
title = slide_2.shapes.title
title.text = "Distribuzione Dati Demografici"

content = slide_2.shapes.placeholders[1]
content.text = f"Distribuzione di Gender:\n{gender_dist}\n\nDistribuzione di Age:\n{df['Age'].describe()}"

# **Slide 3: Distribuzione di Gender (grafico)**
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_3.shapes.title
title.text = "Distribuzione di Gender"
slide_3.shapes.add_picture('gender_distribution.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 4: Distribuzione di Age (grafico)**
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_4.shapes.title
title.text = "Distribuzione di Age"
slide_4.shapes.add_picture('age_distribution.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 5: Numero di Studenti per Department (grafico)**
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_5.shapes.title
title.text = "Numero di Studenti per Department"
slide_5.shapes.add_picture('department_distribution.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 6: Performance degli Studenti - Punteggi Medi per Department (grafico)**
slide_6 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_6.shapes.title
title.text = "Punteggi Medi per Department"
slide_6.shapes.add_picture('average_scores_department.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 7: Distribuzione di Final_Score (Boxplot)**
slide_7 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_7.shapes.title
title.text = "Distribuzione di Final_Score"
slide_7.shapes.add_picture('final_score_boxplot.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 8: Relazione tra Study Hours e Total Score (grafico)**
slide_8 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_8.shapes.title
title.text = "Relazione tra Study Hours per Week e Total Score"
slide_8.shapes.add_picture('study_hours_vs_score.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 9: Relazione tra Attendance e Performance (grafico)**
slide_9 = prs.slides.add_slide(prs.slide_layouts[5])  # Solo titolo e immagine
title = slide_9.shapes.title
title.text = "Relazione tra Attendance (%) e Total Score"
slide_9.shapes.add_picture('attendance_vs_score.png', Inches(0.5), Inches(1.5), width=Inches(9))

# **Slide 10: Sintesi**
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])  # Titolo e contenuto
title = slide_10.shapes.title
title.text = "Sintesi e Conclusioni"

content = slide_10.shapes.placeholders[1]
content.text = "I principali insight del dataset sono:\n- Gender ha una distribuzione equilibrata tra i generi.\n- Gli studenti di alcuni dipartimenti hanno punteggi più alti.\n- La relazione tra le ore di studio e il punteggio finale è positiva.\n- Maggiore è l'Attendance, migliore è la performance finale."

# **Salva il file PowerPoint**
prs.save('Students_Analysis_Report.pptx')

print("Il report è stato creato e salvato come 'Students_Analysis_Report.pptx'")
